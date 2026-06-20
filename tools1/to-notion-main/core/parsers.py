import os
import email
import shutil
import json
import csv
import io
import logging
import zipfile
import tarfile
import tempfile
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
from pdf2image import convert_from_path

logger = logging.getLogger(__name__)

def is_tesseract_available():
    return shutil.which("tesseract") is not None

def ocr_image(file_path):
    if not is_tesseract_available():
        return "[OCR not available: Tesseract not found]"
    try:
        return pytesseract.image_to_string(Image.open(file_path))
    except Exception as e:
        return f"[OCR Error: {e}]"

def parse_pdf(file_path):
    reader = PdfReader(file_path)
    all_pages = []
    current_chunk = []

    # Try text extraction first
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""

        # Heuristic to detect garbage or empty text
        # If less than 10% of characters are alphanumeric, it might be an image-only page or bad extraction
        alnum_count = sum(c.isalnum() for c in page_text)
        is_garbage = len(page_text) > 0 and (alnum_count / len(page_text)) < 0.1

        # If no text extracted or it's garbage and OCR is available, try OCRing the page
        if (not page_text.strip() or is_garbage) and is_tesseract_available():
            try:
                images = convert_from_path(file_path, first_page=index, last_page=index)
                if images:
                    page_text = pytesseract.image_to_string(images[0])
            except Exception:
                pass

        if page_text:
            current_chunk.append(f"--- PDF Page {index} ---\n{page_text}")

        if len(current_chunk) == 10:
            all_pages.append("\n\n".join(current_chunk))
            current_chunk = []

    if current_chunk:
        all_pages.append("\n\n".join(current_chunk))
    return all_pages

def parse_docx(file_path):
    doc = Document(file_path)
    text_content = []
    for p in doc.paragraphs:
        if not p.text.strip(): continue

        style_name = p.style.name if p.style and hasattr(p.style, 'name') else ""

        if style_name.startswith('Heading 1'):
            text_content.append(f"# {p.text}")
        elif style_name.startswith('Heading 2'):
            text_content.append(f"## {p.text}")
        elif style_name.startswith('Heading 3'):
            text_content.append(f"### {p.text}")
        elif style_name.startswith('Heading'):
            text_content.append(f"#### {p.text}")
        elif 'List' in style_name or p._element.xpath('./w:pPr/w:numPr'):
            text_content.append(f"- {p.text}")
        else:
            text_content.append(p.text)

    # Chunking by virtual pages of 500 paragraphs each
    chunk_size = 500
    return ["\n".join(text_content[i:i + chunk_size]) for i in range(0, len(text_content), chunk_size)]

def parse_pptx(file_path):
    prs = Presentation(file_path)
    all_slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        all_slides.append("\n".join(slide_text))

    # Chunking every 5 slides
    chunk_size = 5
    return ["\n\n".join(all_slides[i:i + chunk_size]) for i in range(0, len(all_slides), chunk_size)]

def parse_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    return [clean_html_soup(soup)]

def clean_html_soup(soup):
    for script in soup(["script", "style", "nav", "footer"]):
        script.extract()

    text_content = []

    def _get_inline_markdown(element):
        if isinstance(element, str):
            return element

        parts = []
        for child in element.children:
            parts.append(_get_inline_markdown(child))

        content = "".join(parts)
        if not content.strip() and element.name not in ['br', 'hr']:
            return content

        if element.name in ['b', 'strong']:
            return f"**{content}**"
        if element.name in ['i', 'em']:
            return f"*{content}*"
        if element.name == 'u':
            return f"<u>{content}</u>"
        if element.name in ['s', 'strike', 'del']:
            return f"~~{content}~~"
        if element.name == 'a':
            href = element.get('href')
            return f"[{content}]({href})" if href else content
        if element.name == 'code':
            # Only if not inside pre (handled separately)
            p = element.parent
            while p:
                if p.name == 'pre': return content
                p = p.parent
            return f"`{content}`"
        return content

    # Use a more recursive/structural approach instead of just finding all at top level
    def _process_element(element):
        if not element or element.name in ['script', 'style', 'nav', 'footer']:
            return

        if element.name == 'hr':
            text_content.append("---")
            return

        # Direct content check for certain tags
        if element.name == 'table':
            rows = []
            for tr in element.find_all('tr', recursive=False):
                cells = []
                for td in tr.find_all(['td', 'th'], recursive=False):
                    cells.append(_get_inline_markdown(td).strip().replace("|", "\\|"))
                if cells:
                    rows.append("| " + " | ".join(cells) + " |")

            if rows:
                # Add separator
                num_cols = max(len(r.split('|')) - 2 for r in rows)
                if num_cols > 0:
                    separator = "| " + " | ".join(["---"] * num_cols) + " |"
                    rows.insert(1, separator)
                    text_content.append("\n".join(rows))
            return

        if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'pre', 'code', 'blockquote']:
            if element.name in ['pre', 'code'] and (element.name == 'pre' or element.parent.name != 'pre'):
                text = element.get_text()
                if text.strip():
                    text_content.append(f"```\n{text.strip()}\n```")
                return

            text = _get_inline_markdown(element).strip()
            if not text: return

            if element.name == 'h1': text_content.append(f"# {text}")
            elif element.name == 'h2': text_content.append(f"## {text}")
            elif element.name == 'h3': text_content.append(f"### {text}")
            elif element.name in ['h4', 'h5', 'h6']: text_content.append(f"### {text}")
            elif element.name == 'li': text_content.append(f"- {text}")
            elif element.name == 'blockquote': text_content.append(f"> {text}")
            else: text_content.append(text)
            return

        # If it's a container like div or span, look at children
        for child in element.children:
            if hasattr(child, 'name') and child.name:
                _process_element(child)
            elif isinstance(child, str):
                text = child.strip()
                if text:
                    text_content.append(text)

    # If the soup itself is a snippet, we process its children
    for child in soup.children:
        if hasattr(child, 'name') and child.name:
            _process_element(child)
        elif isinstance(child, str):
            text = child.strip()
            if text:
                text_content.append(text)

    return "\n".join(text_content)

def parse_markdown(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return [f.read()]

def parse_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return [f.read()]

def parse_json(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        pretty_json = json.dumps(data, indent=2)
        return [f"```json\n{pretty_json}\n```"]

def parse_yaml(file_path):
    import yaml
    with open(file_path, 'r', encoding='utf-8') as f:
        data = yaml.safe_load(f)
        pretty_yaml = yaml.dump(data, default_flow_style=False)
        return [f"```yaml\n{pretty_yaml}\n```"]

def parse_xml(file_path):
    from xml.dom import minidom
    with open(file_path, 'r', encoding='utf-8') as f:
        xml_str = f.read()
        try:
            reparsed = minidom.parseString(xml_str)
            pretty_xml = reparsed.toprettyxml(indent="  ")
            return [f"```xml\n{pretty_xml}\n```"]
        except Exception:
            return [f"```xml\n{xml_str}\n```"]

def parse_csv(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        rows = list(reader)
        if not rows:
            return [""]

        # Simple Markdown table conversion
        header = rows[0]
        markdown_table = "| " + " | ".join(header) + " |\n"
        markdown_table += "| " + " | ".join(["---"] * len(header)) + " |\n"
        for row in rows[1:]:
            markdown_table += "| " + " | ".join(row) + " |\n"

        return [markdown_table]

def parse_excel(file_path):
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    all_sheets_content = []
    for sheet in wb.worksheets:
        rows = list(sheet.values)
        if not rows: continue

        content = [f"## Sheet: {sheet.title}"]
        header = [str(cell).replace("|", "\\|") if cell is not None else "" for cell in rows[0]]
        markdown_table = "| " + " | ".join(header) + " |\n"
        markdown_table += "| " + " | ".join(["---"] * len(header)) + " |\n"
        for row in rows[1:]:
            row_vals = [str(cell).replace("|", "\\|") if cell is not None else "" for cell in row]
            markdown_table += "| " + " | ".join(row_vals) + " |\n"
        content.append(markdown_table)
        all_sheets_content.append("\n".join(content))

    return all_sheets_content

def parse_mhtml(file_path):
    with open(file_path, 'rb') as f:
        msg = email.message_from_binary_file(f)

    for part in msg.walk():
        if part.get_content_type() == "text/html":
            try:
                html_payload = part.get_payload(decode=True).decode(part.get_content_charset() or 'utf-8', errors='ignore')
                soup = BeautifulSoup(html_payload, 'html.parser')
                return [clean_html_soup(soup)]
            except Exception as e:
                logger.error(f"MHTML sub-part parsing error: {e}")
                continue
    return [""]

def is_within_directory(directory, target):
    abs_directory = os.path.abspath(directory)
    abs_target = os.path.abspath(target)
    prefix = os.path.commonprefix([abs_directory, abs_target])
    return prefix == abs_directory

def extract_archive(file_path, temp_dir):
    """Extracts zip or tar archives to a temporary directory with Path Traversal protection."""
    if file_path.endswith('.zip'):
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            for member in zip_ref.infolist():
                target_path = os.path.join(temp_dir, member.filename)
                if not is_within_directory(temp_dir, target_path):
                    raise Exception(f"Potential Path Traversal attack detected in zip: {member.filename}")
            zip_ref.extractall(temp_dir)
    elif file_path.endswith(('.tar.gz', '.tgz', '.tar')):
        with tarfile.open(file_path, 'r:*') as tar_ref:
            for member in tar_ref.getmembers():
                target_path = os.path.join(temp_dir, member.name)
                if not is_within_directory(temp_dir, target_path):
                    raise Exception(f"Potential Path Traversal attack detected in tar: {member.name}")
            tar_ref.extractall(temp_dir)

def process_uploaded_document(file_path, extension):
    ext = extension.lower()
    if ext in ['.zip', '.tar.gz', '.tgz', '.tar']:
        all_chunks = []
        with tempfile.TemporaryDirectory() as temp_dir:
            extract_archive(file_path, temp_dir)
            for root, _, files in os.walk(temp_dir):
                for file in sorted(files):
                    sub_file_path = os.path.join(root, file)
                    sub_ext = os.path.splitext(file)[1].lower()
                    try:
                        chunks = process_uploaded_document(sub_file_path, sub_ext)
                        if chunks:
                            all_chunks.append(f"## File: {os.path.relpath(sub_file_path, temp_dir)}")
                            all_chunks.extend(chunks)
                    except Exception:
                        continue
        return all_chunks

    if ext == '.pdf': return parse_pdf(file_path)
    elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
        return [ocr_image(file_path)]
    elif ext in ['.docx', '.doc']: return parse_docx(file_path)
    elif ext == '.pptx': return parse_pptx(file_path)
    elif ext in ['.html', '.htm']: return parse_html(file_path)
    elif ext in ['.md', '.markdown']: return parse_markdown(file_path)
    elif ext == '.txt': return parse_txt(file_path)
    elif ext == '.json': return parse_json(file_path)
    elif ext in ['.yaml', '.yml']: return parse_yaml(file_path)
    elif ext == '.xml': return parse_xml(file_path)
    elif ext == '.csv': return parse_csv(file_path)
    elif ext in ['.xlsx', '.xlsm', '.xltx', '.xltm']: return parse_excel(file_path)
    elif ext == '.mhtml': return parse_mhtml(file_path)
    else: raise ValueError(f"Unsupported document type: {extension}")
